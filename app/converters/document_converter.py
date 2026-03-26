import os
import shutil
import subprocess
import logging
from PIL import Image

Image.MAX_IMAGE_PIXELS = None  # Allow large renders (300 DPI full-page images)

logger = logging.getLogger(__name__)


def document_to_pdf(input_path: str, output_path: str, temp_dir: str) -> None:
    """
    Convert DOCX/DOC to PDF using LibreOffice headless.
    LibreOffice writes the output file next to the input, so we move it after.
    """
    cmd = [
        "soffice",
        "--headless",
        "--convert-to", "pdf",
        "--outdir", temp_dir,
        input_path,
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")

    # LibreOffice names the output after the input file
    base = os.path.splitext(os.path.basename(input_path))[0]
    libreoffice_output = os.path.join(temp_dir, f"{base}.pdf")

    if not os.path.exists(libreoffice_output):
        raise RuntimeError("LibreOffice did not produce an output file")

    shutil.move(libreoffice_output, output_path)
    logger.info("Converted document to PDF: %s", output_path)


def markdown_to_pdf(input_path: str, output_path: str) -> None:
    """Convert Markdown to PDF using pandoc."""
    cmd = [
        "pandoc", input_path,
        "-o", output_path,
        "--pdf-engine=weasyprint",
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        raise RuntimeError(f"pandoc Markdown→PDF failed: {result.stderr}")
    logger.info("Converted Markdown to PDF")


def csv_to_xlsx(input_path: str, output_path: str) -> None:
    """Convert CSV to XLSX using openpyxl."""
    import csv
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    with open(input_path, newline="", encoding="utf-8-sig") as f:
        reader = csv.reader(f)
        for row in reader:
            ws.append(row)
    wb.save(output_path)
    logger.info("Converted CSV to XLSX")


def xlsx_to_csv(input_path: str, output_path: str) -> None:
    """Convert XLSX to CSV using openpyxl."""
    import csv
    from openpyxl import load_workbook

    wb = load_workbook(input_path, read_only=True)
    ws = wb.active
    with open(output_path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        for row in ws.iter_rows(values_only=True):
            writer.writerow(row)
    wb.close()
    logger.info("Converted XLSX to CSV")


def pdf_to_odt(input_path: str, output_path: str, temp_dir: str) -> None:
    """Convert PDF to ODT via image-based DOCX then LibreOffice DOCX→ODT."""
    import uuid as _uuid
    uid = _uuid.uuid4().hex

    # Step 1: PDF → DOCX (image-based, no clipping)
    tmp_docx = os.path.join(temp_dir, f"_odt_{uid}.docx")
    pdf_to_docx(input_path, tmp_docx)

    # Step 2: DOCX → ODT via LibreOffice
    user_profile = os.path.join(temp_dir, f"_lo_profile_{uid}")
    cmd = [
        "soffice",
        "--headless",
        f"-env:UserInstallation=file://{user_profile}",
        "--convert-to", "odt",
        "--outdir", temp_dir,
        tmp_docx,
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)

    if os.path.exists(user_profile):
        shutil.rmtree(user_profile, ignore_errors=True)
    if os.path.exists(tmp_docx):
        os.remove(tmp_docx)

    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice DOCX→ODT failed: {result.stderr}")

    libreoffice_output = os.path.join(temp_dir, f"_odt_{uid}.odt")
    if not os.path.exists(libreoffice_output):
        raise RuntimeError("LibreOffice did not produce an ODT file")

    shutil.move(libreoffice_output, output_path)
    logger.info("Converted PDF to ODT: %s", output_path)


def _make_anchor(r_id, width_emu, height_emu, pic_idx):
    """Build a wp:anchor element to float an image at (0,0) on the page."""
    from lxml import etree

    nsmap = {
        'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing',
        'r':  'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'a':  'http://schemas.openxmlformats.org/drawingml/2006/main',
        'pic': 'http://schemas.openxmlformats.org/drawingml/2006/picture',
        'wp14': 'http://schemas.microsoft.com/office/word/2010/wordprocessingDrawing',
    }

    anchor_xml = f'''
    <wp:anchor distT="0" distB="0" distL="0" distR="0"
               simplePos="1" relativeHeight="{pic_idx}"
               behindDoc="1" locked="0" layoutInCell="1" allowOverlap="1"
               xmlns:wp="{nsmap['wp']}"
               xmlns:r="{nsmap['r']}"
               xmlns:a="{nsmap['a']}"
               xmlns:pic="{nsmap['pic']}"
               xmlns:wp14="{nsmap['wp14']}">
        <wp:simplePos x="0" y="0"/>
        <wp:positionH relativeFrom="page"><wp:posOffset>0</wp:posOffset></wp:positionH>
        <wp:positionV relativeFrom="page"><wp:posOffset>0</wp:posOffset></wp:positionV>
        <wp:extent cx="{width_emu}" cy="{height_emu}"/>
        <wp:effectExtent l="0" t="0" r="0" b="0"/>
        <wp:wrapNone/>
        <wp:docPr id="{pic_idx}" name="Picture {pic_idx}"/>
        <a:graphic>
            <a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/picture">
                <pic:pic>
                    <pic:nvPicPr>
                        <pic:cNvPr id="{pic_idx}" name="img{pic_idx}.png"/>
                        <pic:cNvPicPr/>
                    </pic:nvPicPr>
                    <pic:blipFill>
                        <a:blip r:embed="{r_id}"/>
                        <a:stretch><a:fillRect/></a:stretch>
                    </pic:blipFill>
                    <pic:spPr>
                        <a:xfrm>
                            <a:off x="0" y="0"/>
                            <a:ext cx="{width_emu}" cy="{height_emu}"/>
                        </a:xfrm>
                        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
                    </pic:spPr>
                </pic:pic>
            </a:graphicData>
        </a:graphic>
    </wp:anchor>
    '''
    return etree.fromstring(anchor_xml)


def pdf_to_docx(input_path: str, output_path: str) -> None:
    """
    Convert PDF to DOCX by rendering each page as a high-res image
    and placing it as a floating anchor at (0,0). This avoids clipping
    from paragraph line-height overhead.
    """
    from docx import Document
    from docx.shared import Emu
    from docx.oxml.ns import qn
    from pypdf import PdfReader, PdfWriter

    import uuid as _uuid

    reader = PdfReader(input_path)
    page_count = len(reader.pages)
    doc = Document()
    temp_dir = os.path.dirname(output_path)
    uid = _uuid.uuid4().hex

    # Remove the default empty paragraph
    for p in doc.paragraphs:
        p._element.getparent().remove(p._element)

    img_paths = []
    try:
        for i in range(page_count):
            # Extract single page
            single_pdf = os.path.join(temp_dir, f"_docx_{uid}_p{i}.pdf")
            writer = PdfWriter()
            writer.add_page(reader.pages[i])
            with open(single_pdf, "wb") as f:
                writer.write(f)

            # Render to PNG at 300 DPI via Inkscape
            img_path = os.path.join(temp_dir, f"_docx_{uid}_p{i}.png")
            cmd = [
                "inkscape", "--pdf-poppler", single_pdf,
                "--export-type=png", f"--export-filename={img_path}",
                "--export-dpi=300",
            ]
            result = subprocess.run(cmd, capture_output=True, text=True)
            os.remove(single_pdf)
            if result.returncode != 0:
                raise RuntimeError(f"Inkscape render failed on page {i+1}: {result.stderr}")

            img_paths.append(img_path)

            # Use actual rendered image dimensions
            with Image.open(img_path) as img:
                img_w, img_h = img.size

            # Image aspect ratio
            aspect = img_w / img_h

            # A4 dimensions in EMU (210mm x 297mm)
            A4_SHORT = int(210 / 25.4 * 914400)  # ~7560000
            A4_LONG = int(297 / 25.4 * 914400)   # ~10692000

            # Choose landscape or portrait based on image aspect ratio
            if aspect >= 1.0:
                # Landscape
                section_w = A4_LONG
                section_h = A4_SHORT
            else:
                # Portrait
                section_w = A4_SHORT
                section_h = A4_LONG

            # Scale image to fit within the page, preserving aspect ratio
            img_w_emu = int(img_w / 300 * 914400)
            img_h_emu = int(img_h / 300 * 914400)
            scale = min(section_w / img_w_emu, section_h / img_h_emu)
            fit_w = int(img_w_emu * scale)
            fit_h = int(img_h_emu * scale)

            if i == 0:
                section = doc.sections[0]
            else:
                section = doc.add_section()
            section.page_width = Emu(section_w)
            section.page_height = Emu(section_h)
            section.top_margin = Emu(0)
            section.bottom_margin = Emu(0)
            section.left_margin = Emu(0)
            section.right_margin = Emu(0)
            section.header_distance = Emu(0)
            section.footer_distance = Emu(0)

            # Add image inline with line spacing matching section height
            paragraph = doc.add_paragraph()
            pPr = paragraph._element.get_or_add_pPr()
            line_twips = str(int(section_h / 12700 * 20) + 1)
            spacing = pPr.makeelement(qn('w:spacing'), {
                qn('w:before'): '0',
                qn('w:after'): '0',
                qn('w:line'): line_twips,
                qn('w:lineRule'): 'exact',
            })
            pPr.append(spacing)
            run = paragraph.add_run()
            run.add_picture(img_path, width=Emu(fit_w), height=Emu(fit_h))

            logger.info("Added page %d/%d to DOCX", i + 1, page_count)

        doc.save(output_path)
        logger.info("Converted PDF to DOCX: %s", output_path)
    finally:
        for p in img_paths:
            if os.path.exists(p):
                os.remove(p)


def pdf_to_pptx(input_path: str, output_path: str, temp_dir: str) -> None:
    """Convert PDF to PPTX via SVG→EMF. Each page becomes one vector object per slide."""
    import uuid as _uuid
    from pptx import Presentation
    from pptx.util import Emu
    from pypdf import PdfReader, PdfWriter

    reader = PdfReader(input_path)
    page_count = len(reader.pages)
    prs = Presentation()
    uid = _uuid.uuid4().hex

    tmp_files = []
    try:
        for i in range(page_count):
            # Extract single page
            single_pdf = os.path.join(temp_dir, f"_pptx_{uid}_p{i}.pdf")
            writer = PdfWriter()
            writer.add_page(reader.pages[i])
            with open(single_pdf, "wb") as f:
                writer.write(f)
            tmp_files.append(single_pdf)

            # Get page dimensions from PDF (in points, 1pt = 1/72 inch)
            page = reader.pages[i]
            pdf_w = float(page.mediabox.width)
            pdf_h = float(page.mediabox.height)

            # Convert PDF → SVG via Inkscape
            svg_path = os.path.join(temp_dir, f"_pptx_{uid}_p{i}.svg")
            cmd_svg = [
                "inkscape", "--pdf-poppler", single_pdf,
                "--export-type=svg", f"--export-filename={svg_path}",
            ]
            result = subprocess.run(cmd_svg, capture_output=True, text=True)
            if result.returncode != 0:
                raise RuntimeError(f"Inkscape PDF→SVG failed on page {i+1}: {result.stderr}")
            tmp_files.append(svg_path)

            # Convert SVG → EMF via Inkscape (vector format supported by PowerPoint)
            emf_path = os.path.join(temp_dir, f"_pptx_{uid}_p{i}.emf")
            cmd_emf = [
                "inkscape", svg_path,
                "--export-type=emf", f"--export-filename={emf_path}",
            ]
            result = subprocess.run(cmd_emf, capture_output=True, text=True)
            if result.returncode != 0:
                raise RuntimeError(f"Inkscape SVG→EMF failed on page {i+1}: {result.stderr}")
            tmp_files.append(emf_path)

            # Set slide size from first page (points → EMU: 1pt = 12700 EMU)
            slide_w = Emu(int(pdf_w * 12700))
            slide_h = Emu(int(pdf_h * 12700))

            if i == 0:
                prs.slide_width = slide_w
                prs.slide_height = slide_h

            # Add blank slide with EMF as single object
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(emf_path, 0, 0, slide_w, slide_h)

            logger.info("Added page %d/%d to PPTX (vector)", i + 1, page_count)

        prs.save(output_path)
        logger.info("Converted PDF to PPTX: %s", output_path)
    finally:
        for p in tmp_files:
            if os.path.exists(p):
                os.remove(p)
